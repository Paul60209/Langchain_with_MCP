# Fix compatibility issue with python-pptx in Python 3.10+ first
import collections
if not hasattr(collections, 'Container'):
    import collections.abc
    collections.Container = collections.abc.Container
    collections.Mapping = collections.abc.Mapping
    collections.MutableMapping = collections.abc.MutableMapping
    collections.Sequence = collections.abc.Sequence
    collections.Set = collections.abc.Set

# Now import required modules
import tempfile
import os
import dotenv
import base64
import argparse
import json

# Import OpenAI API
from langchain_openai import ChatOpenAI

# Import MCP related modules
from mcp.server.fastmcp import FastMCP

# Finally, import python-pptx related modules
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor

# Load environment variables
dotenv.load_dotenv()

# Define output path
OUTPUT_PATH = 'output'

# Create MCP server
mcp = FastMCP("PPTTranslatorServer")

# Parse command line arguments
parser = argparse.ArgumentParser(description='PPT Translator MCP Server')
parser.add_argument('--port', type=int, default=8003, help='Server listening port (default: 8003)')
args = parser.parse_args()

# Set environment variable for FastMCP to use the specified port
os.environ["MCP_SSE_PORT"] = str(args.port)

def get_text_frame_properties(text_frame):
    """Get all formatting properties of a text frame."""
    properties = {
        'margin_left': text_frame.margin_left,
        'margin_right': text_frame.margin_right,
        'margin_top': text_frame.margin_top,
        'margin_bottom': text_frame.margin_bottom,
        'vertical_anchor': text_frame.vertical_anchor,
        'word_wrap': text_frame.word_wrap,
        'auto_size': text_frame.auto_size,
    }
    return properties

def get_paragraph_properties(paragraph):
    """Get all formatting properties of a paragraph."""
    properties = {
        'alignment': paragraph.alignment,
        'level': paragraph.level,
        'line_spacing': paragraph.line_spacing,
        'space_before': paragraph.space_before,
        'space_after': paragraph.space_after,
    }
    return properties

def get_color_properties(color):
    """Get color properties."""
    if not color:
        return None

    properties = {
        'type': color.type if hasattr(color, 'type') else None,
        'rgb': color.rgb if hasattr(color, 'rgb') else None,
        'theme_color': color.theme_color if hasattr(color, 'theme_color') else None,
        'brightness': color.brightness if hasattr(color, 'brightness') else None,
    }
    return properties

def get_run_properties(run):
    """Get all formatting properties of a text run."""
    font = run.font
    properties = {
        'size': font.size,
        'name': font.name,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color': get_color_properties(font.color),
        'fill': get_color_properties(font.fill.fore_color) if hasattr(font, 'fill') else None,
    }
    return properties

def apply_color_properties(color_obj, properties):
    """Apply color properties."""
    if not properties or not color_obj:
        return

    try:
        # If RGB value exists, set RGB color directly
        if properties['rgb']:
            if isinstance(properties['rgb'], (tuple, list)) and len(properties['rgb']) == 3:
                color_obj.rgb = RGBColor(*properties['rgb'])
            else:
                color_obj.rgb = properties['rgb']
        # If theme color exists, set theme color
        elif properties['theme_color'] and properties['theme_color'] != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
            color_obj.theme_color = properties['theme_color']
            if properties['brightness'] is not None:
                color_obj.brightness = properties['brightness']
    except Exception as e:
        print(f"Error setting color: {str(e)}")
        pass  # If setting fails, keep the original color

def apply_text_frame_properties(text_frame, properties):
    """Apply text frame formatting properties."""
    text_frame.margin_left = properties['margin_left']
    text_frame.margin_right = properties['margin_right']
    text_frame.margin_top = properties['margin_top']
    text_frame.margin_bottom = properties['margin_bottom']
    text_frame.vertical_anchor = properties['vertical_anchor']
    text_frame.word_wrap = properties['word_wrap']
    text_frame.auto_size = properties['auto_size']

def apply_paragraph_properties(paragraph, properties):
    """Apply paragraph formatting properties."""
    paragraph.alignment = properties['alignment']
    paragraph.level = properties['level']
    paragraph.line_spacing = properties['line_spacing']
    paragraph.space_before = properties['space_before']
    paragraph.space_after = properties['space_after']

def apply_run_properties(run, properties):
    """Apply text run formatting properties."""
    font = run.font
    if properties['size']:
        font.size = properties['size']
    if properties['name']:
        font.name = properties['name']
    if properties['bold'] is not None:
        font.bold = properties['bold']
    if properties['italic'] is not None:
        font.italic = properties['italic']
    if properties['underline'] is not None:
        font.underline = properties['underline']

    # Apply color
    if properties['color']:
        apply_color_properties(font.color, properties['color'])
    if properties['fill'] and hasattr(font, 'fill'):
        apply_color_properties(font.fill.fore_color, properties['fill'])

async def translate_text(text: str, olang: str, tlang: str, ctx=None) -> str:
    """Translate text using ChatGPT.

    Args:
        text (str): Text to be translated
        olang (str): Original language code
        tlang (str): Target language code
        ctx: MCP context object (no longer used)

    Returns:
        str: Translated text
    """
    if not text.strip():
        return text

    print(f"\nTranslating text:")
    print(f"Original ({olang}): {text}")

    try:
        # Create ChatGPT model
        model = ChatOpenAI(temperature=0)

        # Create system prompt
        system_message = f"""You are a professional translator. Translate the following text from {olang} to {tlang}.
        Rules:
        1. Keep all formatting symbols (like bullet points, numbers) unchanged
        2. Keep all special characters unchanged
        3. Keep all whitespace and line breaks
        4. Only translate the actual text content
        5. Maintain the same tone and style
        6. Do not add any explanations or notes
        7. Keep all numbers and dates unchanged
        8. Keep all proper nouns unchanged unless they have standard translations
        """

        # Create message list
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": text}
        ]

        # Execute translation
        response = await model.ainvoke(messages)
        translated_text = response.content.strip()

        print(f"Translation ({tlang}): {translated_text}\n")
        return translated_text

    except Exception as e:
        print(f"Translation failed: {str(e)}")
        if ctx:
            await ctx.info(f"Translation failed: {str(e)}")
        # Return original text to ensure content is not lost
        return text

async def translate_group_shape(shape, olang: str, tlang: str, ctx=None) -> None:
    """Translate all shapes within a group.

    Args:
        shape: PowerPoint group shape object
        olang (str): Original language code
        tlang (str): Target language code
        ctx: MCP context object
    """
    try:
        if not hasattr(shape, 'shapes'):
            return

        # Iterate through all shapes in the group
        for child_shape in shape.shapes:
            if child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recursively handle nested groups
                await translate_group_shape(child_shape, olang, tlang, ctx)
            else:
                # Translate individual shape
                await translate_shape(child_shape, olang, tlang, ctx)
    except Exception as e:
        print(f"Error translating group shape: {str(e)}")
        if ctx:
            await ctx.info(f"Error translating group shape: {str(e)}")
        raise

async def translate_shape(shape, olang: str, tlang: str, ctx=None) -> None:
    """Translate a shape in PowerPoint.

    Args:
        shape: PowerPoint shape object
        olang (str): Original language code
        tlang (str): Target language code
        ctx: MCP context object
    """
    try:
        # Handle group shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:

            await translate_group_shape(shape, olang, tlang, ctx)
            return

        # Check if the shape contains a text frame
        if not hasattr(shape, "text_frame"):

            return

        text_frame = shape.text_frame
        if not text_frame.text.strip():

            return

        # Save text frame format
        text_frame_props = get_text_frame_properties(text_frame)

        # Iterate through all paragraphs
        paragraph_count = len(text_frame.paragraphs)

        for i, paragraph in enumerate(text_frame.paragraphs, 1):
            # Save paragraph format
            para_props = get_paragraph_properties(paragraph)

            # Iterate through all text runs
            runs_data = []
            run_count = len(paragraph.runs)

            for j, run in enumerate(paragraph.runs, 1):
                # Save run format and text
                run_props = get_run_properties(run)
                original_text = run.text

                if original_text.strip():
                    # Translate using ChatGPT
                    translated_text = await translate_text(original_text, olang, tlang, ctx)
                    runs_data.append((translated_text, run_props))
                else:
                    # Skip empty run text
                    runs_data.append((original_text, run_props))

            # Clear original content
            for _ in range(len(paragraph.runs)):
                paragraph._p.remove(paragraph.runs[0]._r)

            # Add translated text and apply format
            for text, props in runs_data:
                run = paragraph.add_run()
                run.text = text
                apply_run_properties(run, props)

            # Restore paragraph format
            apply_paragraph_properties(paragraph, para_props)

        # Restore text frame format
        apply_text_frame_properties(text_frame, text_frame_props)

    except Exception as e:
        print(f"[Error] Error translating shape: {str(e)}")
        if ctx:
            await ctx.info(f"Error translating shape: {str(e)}")
        raise

async def translate_ppt_file(file_path: str, olang: str, tlang: str, ctx=None) -> str:
    """Translate a PowerPoint file.

    Args:
        file_path (str): Path to the PowerPoint file
        olang (str): Original language code
        tlang (str): Target language code
        ctx: MCP context object

    Returns:
        str: Path to the translated file
    """
    try:
        print("\n========== PowerPoint Translation Process ==========")
        # 1. Create output directory
        os.makedirs(OUTPUT_PATH, exist_ok=True)

        # 2. Prepare output file path
        file_name = os.path.basename(file_path)
        name, ext = os.path.splitext(file_name)
        output_file = f'translated_{name}{ext}'
        output_path = os.path.join(OUTPUT_PATH, output_file)

        # 3. Load PowerPoint
        print(f"[Info] Starting PowerPoint translation...")
        print(f"[Info] Source language: {olang}")
        print(f"[Info] Target language: {tlang}")
        if ctx:
            await ctx.info(f"Starting translation...\nFrom {olang} to {tlang}")

        presentation = Presentation(file_path)
        total_slides = len(presentation.slides)

        # 4. Translate each slide
        for index, slide in enumerate(presentation.slides, 1):
            progress_msg = f"Translating slide {index}/{total_slides}..."
            print(f"\n[Progress] {progress_msg}")

            # Count shapes in the slide
            shape_count = len(slide.shapes)

            if ctx:
                await ctx.info(progress_msg)
                # Report progress (0-100%)
                await ctx.report_progress(index - 1, total_slides)

            # Translate each shape in the slide
            for shape_idx, shape in enumerate(slide.shapes, 1):

                await translate_shape(shape, olang, tlang, ctx)

        # 5. Save the translated file
        if ctx:
            await ctx.info("Translation complete, generating file...")
            await ctx.report_progress(total_slides, total_slides)  # 100% complete

        presentation.save(output_path)
        print(f"[Success] Translated file saved to: {output_path}")
        print(f"========== PowerPoint Translation Complete ==========\n")

        # 6. Return the path
        return output_path

    except Exception as e:
        error_msg = f"Error during translation process: {str(e)}"
        print(f"\n[Error] {error_msg}")
        if ctx:
            await ctx.info(error_msg)
        raise

@mcp.tool()
async def translate_ppt(olang: str, tlang: str, file_content: str = None, file_name: str = None) -> str:
    """
    Translate a PowerPoint file from one language to another while preserving the original format.

    ## Usage Scenario
    - When needing to translate presentations into other languages
    - Preparing multilingual presentations
    - International conference/speech preparation

    ## Parameter Description
    :param olang: Source language code or name, e.g., 'zh-TW', 'Traditional Chinese', 'english', etc.
    :param tlang: Target language code or name, e.g., 'en', 'English', 'japanese', etc.
    :param file_content: Content of the PowerPoint file (base64 encoded string)
    :param file_name: File name (optional, used to determine file type)

    ## Input Example
    - From Chinese to English: olang="Chinese", tlang="English"
    - From English to Japanese: olang="english", tlang="japanese"
    - From Japanese to Chinese: olang="ja", tlang="zh-TW"

    ## File Requirements
    - Supports .ppt and .pptx formats
    - File size should not exceed 10MB
    - Preserves original formatting including font, color, layout, etc.

    ## Notes
    - Translating large files may take several minutes
    - Complex charts and special formats might not be perfectly preserved
    - Proper nouns may require manual correction

    :return: JSON string containing the translation result message and file content
    """
    # Get MCP context
    try:
        ctx = mcp.get_current_request_context()

    except Exception as e:
        print(f"[Warning] Could not get MCP context: {str(e)}")
        ctx = None

    try:
        print(f"\n========== PPT Translator Tool Started ==========")
        print(f"[Parameter] Source language: {olang}")
        print(f"[Parameter] Target language: {tlang}")



        # Check necessary parameters
        if not file_content:
            print(f"[Error] File content not provided")
            return json.dumps({
                "success": False,
                "message": "Error: PowerPoint file content is required. Please upload the file and provide the content."
            })

        # Ensure file name is valid
        if not file_name:
            file_name = "uploaded_presentation.pptx"

        elif not (file_name.lower().endswith('.ppt') or file_name.lower().endswith('.pptx')):
            old_name = file_name
            file_name += ".pptx"  # Add default extension

        # Create a temporary file to store the uploaded content
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as temp_file:
            # Write file content directly
            try:
                if isinstance(file_content, bytes):

                    temp_file.write(file_content)
                elif isinstance(file_content, str):

                    # Try to decode base64 string
                    try:
                        decoded_content = base64.b64decode(file_content)
                        temp_file.write(decoded_content)

                    except Exception as e:
                        print(f"[Warning] Base64 decoding failed: {str(e)}, attempting to write plain text")
                        # If not valid base64, treat as plain text
                        encoded_content = file_content.encode('utf-8')
                        temp_file.write(encoded_content)

                else:
                    print(f"[Error] Unsupported file content type: {type(file_content)}")
                    return json.dumps({
                        "success": False,
                        "message": "Error: Unsupported file content format. Please provide binary or base64 encoded file content."
                    })

                temp_file_path = temp_file.name

            except Exception as e:
                print(f"[Error] Error processing file content: {str(e)}")
                return json.dumps({
                    "success": False,
                    "message": f"Error processing file content: {str(e)}."
                })

        # Execute translation - directly use user-provided language parameters without validation or conversion
        print("[Info] Starting translation process...")
        output_path = await translate_ppt_file(temp_file_path, olang, tlang, ctx)
        print(f"[Info] Translation complete, result path: {output_path}")

        # Clean up temporary file
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)

        # Read the translated file and encode as base64
        with open(output_path, "rb") as f:
            file_bytes = f.read()

            translated_file_content = base64.b64encode(file_bytes).decode('utf-8')

        # Get the output file name
        output_file_name = os.path.basename(output_path)

        # Return JSON containing necessary information
        result_json = json.dumps({
            "success": True,
            "message": "Translation complete!",
            "file_name": output_file_name,
            "file_content": translated_file_content
        })

        print(f"========== PPT Translator Tool Finished ==========\n")
        return result_json

    except Exception as e:
        print(f"[Error] Translator tool execution error: {str(e)}")
        error_result = json.dumps({
            "success": False,
            "message": f"Error during translation process: {str(e)}"
        })
        print(f"========== PPT Translator Tool Error ==========\n")
        return error_result

@mcp.resource("translator://instructions")
async def get_instructions() -> str:
    """Get usage instructions for the PPT translator."""
    return 
    """
    # PowerPoint Translator Tool Usage Instructions

    This tool can translate PowerPoint files from one language to another while preserving the original format.

    ## Supported Languages

    - Chinese (zh-TW)
    - English (en)
    - Japanese (ja)
    - Other language codes can also be attempted

    ## How to Use

    1. Convert the PowerPoint file (.ppt or .pptx) to base64 encoding
    2. Call the `translate_ppt` tool, providing the following parameters:
    - olang: Original language code
    - tlang: Target language code
    - file_content: Base64 encoded content of the file

    ## Translation Process

    1. The tool parses all text in the PowerPoint file
    2. Translates each text element one by one, preserving original format
    3. Generates a new PowerPoint file and returns the result

    ## Notes

    - Translating large files may take longer
    - Some complex formats may not be perfectly preserved
    - File size limit is 10MB
    """

# Start server
if __name__ == "__main__":
    # Ensure output directory exists
    os.makedirs(OUTPUT_PATH, exist_ok=True)

    import uvicorn
    from starlette.applications import Starlette
    from starlette.routing import Route, Mount
    from mcp.server.sse import SseServerTransport

    # Start MCP server
    print(f"Starting PPT Translator Server on port {args.port}")

    # Create SSE transport
    sse = SseServerTransport("/mcp/")

    # Define SSE connection handler function
    async def handle_sse(request):
        async with sse.connect_sse(
            request.scope, request.receive, request._send
        ) as streams:
            await mcp._mcp_server.run(
                streams[0],
                streams[1],
                mcp._mcp_server.create_initialization_options()
            )

    # Create Starlette application
    starlette_app = Starlette(
        routes=[
            # SSE endpoint
            Route("/sse", endpoint=handle_sse, methods=["GET"]),
            Mount("/mcp/", app=sse.handle_post_message)
        ]
    )

    # Start the application using uvicorn
    uvicorn.run(starlette_app, host="0.0.0.0", port=args.port) 